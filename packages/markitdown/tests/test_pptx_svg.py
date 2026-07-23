#!/usr/bin/env python3 -m pytest
"""Tests for PPTX SVG images that lack a rasterized fallback.

PowerPoint stores an SVG picture as an ``<a:blip>`` whose ``r:embed`` points to
a rasterized PNG fallback, plus an ``<asvg:svgBlip>`` extension that points to
the SVG. When a picture has no raster fallback the ``<a:blip>`` has no
``r:embed`` at all, so python-pptx's ``shape.image`` raises
``ValueError("no embedded image")``. The converter must handle this gracefully
(resolving the SVG blip directly) instead of failing the whole conversion.
"""
import os

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from markitdown import MarkItDown
from markitdown.converters._pptx_converter import PptxConverter

TEST_FILES_DIR = os.path.join(os.path.dirname(__file__), "test_files")

# A tiny synthetic PPTX whose only picture is an SVG without a rasterized
# fallback: the <a:blip> has no r:embed, only an <asvg:svgBlip> extension
# pointing to an embedded SVG. Its alt text is "Red square SVG".
SVG_NO_FALLBACK_PPTX = os.path.join(TEST_FILES_DIR, "test_svg_no_fallback.pptx")

_SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


def _first_picture_shape(pptx_path):
    presentation = Presentation(pptx_path)
    converter = PptxConverter()
    for slide in presentation.slides:
        for shape in slide.shapes:
            if converter._is_picture(shape):
                return converter, shape
    raise AssertionError(f"No picture shape found in {pptx_path}")


def test_pptx_svg_without_raster_fallback() -> None:
    md = MarkItDown()

    # Default conversion should not raise and should emit the image alt text.
    result = md.convert(SVG_NO_FALLBACK_PPTX)
    assert "Red square SVG" in result.markdown

    # keep_data_uris used to crash with ValueError("no embedded image"). It
    # should now embed the SVG as a data URI.
    result = md.convert(SVG_NO_FALLBACK_PPTX, keep_data_uris=True)
    assert "data:image/svg+xml;base64," in result.markdown


def test_get_image_info_resolves_svg_blip_without_fallback() -> None:
    # Unit-level check: _get_image_info must resolve the raw SVG blob directly
    # from the <asvg:svgBlip> extension when shape.image raises because there is
    # no rasterized fallback (no r:embed on the <a:blip>).
    converter, shape = _first_picture_shape(SVG_NO_FALLBACK_PPTX)

    blob, content_type, _filename = converter._get_image_info(shape)

    assert blob is not None and len(blob) > 0
    assert content_type == "image/svg+xml"
    assert b"<svg" in blob[:512].lower()


class _FakePart:
    """Minimal part whose related_part always resolves to a truthy object."""

    def related_part(self, rid):
        return object()


class _FakeSvgPlaceholderShape:
    """A placeholder shape whose ``image`` raises like an SVG-only placeholder.

    ``_is_picture`` used to rely on ``hasattr(shape, "image")`` which only
    swallows ``AttributeError`` and let this ``ValueError`` propagate, failing
    even the default conversion.
    """

    shape_type = MSO_SHAPE_TYPE.PLACEHOLDER

    def __init__(self):
        xml = (
            '<p:pic xmlns:p="http://schemas.openxmlformats.org/'
            'presentationml/2006/main" xmlns:a="http://schemas.'
            'openxmlformats.org/drawingml/2006/main" xmlns:r="http://'
            'schemas.openxmlformats.org/officeDocument/2006/relationships" '
            'xmlns:asvg="%s"><a:blip><a:extLst><a:ext uri="{96DAC541-7B7A-'
            '43D3-8B79-37D633B846F1}"><asvg:svgBlip r:embed="rId9"/></a:ext>'
            "</a:extLst></a:blip></p:pic>" % _SVG_NS
        )
        self._element = etree.fromstring(xml)
        self.part = _FakePart()

    @property
    def image(self):
        raise ValueError("no embedded image")


def test_is_picture_true_for_svg_placeholder() -> None:
    # _is_picture must not let shape.image's ValueError propagate for SVG
    # placeholders lacking a raster fallback; it should still report a picture
    # because an embedded SVG blip is present.
    converter = PptxConverter()
    assert converter._is_picture(_FakeSvgPlaceholderShape()) is True


if __name__ == "__main__":
    test_pptx_svg_without_raster_fallback()
    test_get_image_info_resolves_svg_blip_without_fallback()
    test_is_picture_true_for_svg_placeholder()
    print("All tests passed!")
